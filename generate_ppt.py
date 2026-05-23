import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize Presentation and set Widescreen 16:9
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # Color Palette & Styling Constants
    # ----------------------------------------------------
    COLOR_NAVY = RGBColor(15, 23, 42)      # #0F172A (Deep Slate/Navy)
    COLOR_CYAN = RGBColor(14, 165, 233)    # #0EA5E9 (Electric Cyan)
    COLOR_WHITE = RGBColor(255, 255, 255)  # Pure White
    COLOR_LIGHT_BG = RGBColor(248, 250, 252) # #F8FAFC (Slate 50)
    COLOR_TEXT_DARK = RGBColor(30, 41, 59)  # #1E293B (Slate 800)
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # #64748B (Slate 500)
    COLOR_BORDER = RGBColor(226, 232, 240) # #E2E8F0 (Slate 200)
    
    FONT_TITLE = "Segoe UI"
    FONT_BODY = "Arial"

    # ----------------------------------------------------
    # Helper: Set Solid Background Color
    # ----------------------------------------------------
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # ----------------------------------------------------
    # Helper: Add Styled Header to Content Slides
    # ----------------------------------------------------
    def add_slide_header(slide, category, title_text):
        # Category/Tracker Tag (e.g., "SYSTEM ARCHITECTURE")
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = FONT_TITLE
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN
        
        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_NAVY
        
        # Thin Divider Line
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.03)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BORDER
        shape.line.fill.background()

    # ----------------------------------------------------
    # Slide 1: Title Slide (Dark Theme)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1, COLOR_NAVY)
    
    # Left accent block
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_CYAN
    accent_bar.line.fill.background()
    
    # Text container
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(4.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p_tag = tf1.paragraphs[0]
    p_tag.text = "DATA SCIENCE PROJECT SHOWCASE"
    p_tag.font.name = FONT_TITLE
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_CYAN
    p_tag.space_after = Pt(14)
    
    p_title = tf1.add_paragraph()
    p_title.text = "RAG-Based AI Teaching Assistant"
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(46)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.space_after = Pt(10)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "Local Video-to-Text Transcription, Semantic Search, and Q&A System"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_WHITE
    p_sub.space_after = Pt(40)
    
    p_course = tf1.add_paragraph()
    p_course.text = "Inspired by the Ultimate Job Ready Data Science Course (Code With Harry Project 3)"
    p_course.font.name = FONT_BODY
    p_course.font.size = Pt(13)
    p_course.font.color.rgb = COLOR_CYAN
    p_course.font.italic = True

    # ----------------------------------------------------
    # Slide 2: Problem Statement & Challenge
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2, COLOR_LIGHT_BG)
    add_slide_header(slide2, "Context & Challenge", "The Problem: Navigating Video Lectures")
    
    # Left Content Column (Bullet points)
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(6.5), Inches(5.0))
    tf2 = left_box.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "• Inefficiencies in Modern E-Learning"
    p.font.name = FONT_BODY
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.space_after = Pt(6)
    
    p_sub = tf2.add_paragraph()
    p_sub.text = "Students waste hours scrubbing timelines to find exact explanations of coding topics, definitions, or code configurations in multi-hour video courses."
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sub.space_after = Pt(16)
    
    p2 = tf2.add_paragraph()
    p2.text = "• Limitations of Standard Video Search"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_DARK
    p2.space_after = Pt(6)
    
    p2_sub = tf2.add_paragraph()
    p2_sub.text = "Video search features index only video titles, tags, and standard descriptions. The rich pedagogical content spoken within the actual video audio remains completely hidden."
    p2_sub.font.name = FONT_BODY
    p2_sub.font.size = Pt(14)
    p2_sub.font.color.rgb = COLOR_TEXT_MUTED
    p2_sub.space_after = Pt(16)

    p3 = tf2.add_paragraph()
    p3.text = "• Multilingual and Slang Complexities"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_TEXT_DARK
    p3.space_after = Pt(6)
    
    p3_sub = tf2.add_paragraph()
    p3_sub.text = "Popular courses (like CodeWithHarry) use Hinglish (mix of Hindi & English). Traditional transcription tools fail to accurately parse and index this mixed vernacular."
    p3_sub.font.name = FONT_BODY
    p3_sub.font.size = Pt(14)
    p3_sub.font.color.rgb = COLOR_TEXT_MUTED

    # Right Column: Visual Highlight Box
    right_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.9), Inches(4.733), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = COLOR_NAVY
    right_bg.line.color.rgb = COLOR_CYAN
    right_bg.line.width = Pt(1.5)
    
    tb_right = slide2.shapes.add_textbox(Inches(8.1), Inches(2.2), Inches(4.1), Inches(4.2))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    
    pr1 = tf_right.paragraphs[0]
    pr1.text = "THE OBJECTIVE"
    pr1.font.name = FONT_TITLE
    pr1.font.size = Pt(13)
    pr1.font.bold = True
    pr1.font.color.rgb = COLOR_CYAN
    pr1.space_after = Pt(16)
    
    pr2 = tf_right.add_paragraph()
    pr2.text = "Create an intelligent, local, and cost-effective Retrieval-Augmented Generation (RAG) system that:"
    pr2.font.name = FONT_BODY
    pr2.font.size = Pt(15)
    pr2.font.color.rgb = COLOR_WHITE
    pr2.space_after = Pt(16)
    
    bullets = [
        "Automatically converts video playlists to text.",
        "Translates bilingual instruction (Hindi -> English).",
        "Indexes text chunks semantically with vectors.",
        "Gives natural answers anchored to exact timestamps."
    ]
    for b in bullets:
        pr_b = tf_right.add_paragraph()
        pr_b.text = f"✔  {b}"
        pr_b.font.name = FONT_BODY
        pr_b.font.size = Pt(13)
        pr_b.font.color.rgb = COLOR_WHITE
        pr_b.space_after = Pt(10)

    # ----------------------------------------------------
    # Slide 3: Introduction to RAG & Solution
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3, COLOR_LIGHT_BG)
    add_slide_header(slide3, "System Solution", "Introducing Retrieval-Augmented Generation")
    
    # Left Column (Concept)
    left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.8))
    tf3 = left_box.text_frame
    tf3.word_wrap = True
    
    p = tf3.paragraphs[0]
    p.text = "What is RAG in Education?"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    p_desc = tf3.add_paragraph()
    p_desc.text = "Retrieval-Augmented Generation (RAG) integrates semantic search with large language models. Instead of relying purely on a model's static pre-trained weights (which lead to hallucinations), RAG queries a reliable database of course transcripts, extracts relevant context, and passes it to the LLM to generate precise, grounded answers."
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = COLOR_TEXT_DARK
    p_desc.space_after = Pt(20)

    p_val = tf3.add_paragraph()
    p_val.text = "Key Benefits for Students:"
    p_val.font.name = FONT_TITLE
    p_val.font.size = Pt(16)
    p_val.font.bold = True
    p_val.font.color.rgb = COLOR_NAVY
    p_val.space_after = Pt(10)
    
    student_benefits = [
        "Chat directly with the lecture materials.",
        "Zero hallucination answers: derived strictly from the videos.",
        "Hyperlinked direct references to precise seconds in the timeline."
    ]
    for sb in student_benefits:
        psb = tf3.add_paragraph()
        psb.text = f"•  {sb}"
        psb.font.name = FONT_BODY
        psb.font.size = Pt(14)
        psb.font.color.rgb = COLOR_TEXT_MUTED
        psb.space_after = Pt(6)

    # Right Column: Visual Core Pipeline Flow Cards
    cards = [
        ("1. Transcription & Translation", "OpenAI Whisper transcribes multi-lingual audios and translates mixed speech into English structured transcripts."),
        ("2. Vector Embeddings", "Ollama with local model BGE-M3 encodes the textual transcripts into numerical embeddings representing semantic meaning."),
        ("3. Context-Guided LLM Q&A", "Cosine similarity extracts the top relevant sections, and a local LLM (Llama 3.2) generates clear answers with timestamps.")
    ]
    
    for i, (title, desc) in enumerate(cards):
        top_pos = 1.9 + (i * 1.6)
        card_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(top_pos), Inches(5.5), Inches(1.4))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLOR_WHITE
        card_bg.line.color.rgb = COLOR_BORDER
        card_bg.line.width = Pt(1)
        
        card_tb = slide3.shapes.add_textbox(Inches(7.2), Inches(top_pos + 0.1), Inches(5.1), Inches(1.2))
        card_tf = card_tb.text_frame
        card_tf.word_wrap = True
        card_tf.margin_left = card_tf.margin_top = card_tf.margin_bottom = card_tf.margin_right = 0
        
        pct = card_tf.paragraphs[0]
        pct.text = title
        pct.font.name = FONT_TITLE
        pct.font.size = Pt(14)
        pct.font.bold = True
        pct.font.color.rgb = COLOR_NAVY
        pct.space_after = Pt(4)
        
        pcd = card_tf.add_paragraph()
        pcd.text = desc
        pcd.font.name = FONT_BODY
        pcd.font.size = Pt(11.5)
        pcd.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------
    # Slide 4: System Architecture
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4, COLOR_LIGHT_BG)
    add_slide_header(slide4, "Implementation Design", "Core System Architecture & Data Flow")
    
    # 4 horizontal steps block
    steps = [
        ("INGESTION", "video_to_mp3.py", "Extracts audio tracks from MP4 video lectures using bulk subprocess operations and ffmpeg.", "INPUT: .mp4 \nOUTPUT: .mp3"),
        ("TRANSCRIPTION", "mp3_to_json.py", "Utilizes local Whisper large-v2. Translates Hindi to English, structuring segments with start/end timestamps.", "INPUT: .mp3 \nOUTPUT: .json"),
        ("INDEXING", "preprocess_json.py", "Sends text chunks to local Ollama. Generates BGE-M3 embeddings, saving a consolidated joblib DataFrame.", "INPUT: .json \nOUTPUT: .joblib"),
        ("RAG ENGINE", "process_incoming.py", "Accepts student queries, computes Cosine Similarity, extracts Top 5 hits, and prompts local Llama 3.2.", "INPUT: Query \nOUTPUT: Response")
    ]
    
    width = Inches(2.7)
    gap = Inches(0.3)
    start_left = Inches(0.8)
    
    for i, (stage, file_name, desc, io) in enumerate(steps):
        left_pos = start_left + i * (width + gap)
        
        # Step Card
        card = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(2.2), width, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_CYAN if i == 3 else COLOR_BORDER
        card.line.width = Pt(1.5 if i == 3 else 1)
        
        # Text Frame
        tb = slide4.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.35), width - Inches(0.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p1 = tf.paragraphs[0]
        p1.text = f"STAGE 0{i+1}"
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CYAN
        p1.space_after = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = stage
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_NAVY
        p2.space_after = Pt(4)
        
        p3 = tf.add_paragraph()
        p3.text = f"Code: {file_name}"
        p3.font.name = FONT_BODY
        p3.font.size = Pt(10.5)
        p3.font.bold = True
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.space_after = Pt(12)
        
        p4 = tf.add_paragraph()
        p4.text = desc
        p4.font.name = FONT_BODY
        p4.font.size = Pt(11)
        p4.font.color.rgb = COLOR_TEXT_DARK
        p4.space_after = Pt(20)
        
        p5 = tf.add_paragraph()
        p5.text = io
        p5.font.name = FONT_BODY
        p5.font.size = Pt(9.5)
        p5.font.color.rgb = COLOR_CYAN
        p5.font.bold = True

    # ----------------------------------------------------
    # Slide 5: Speech-to-Text Component (Whisper)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5, COLOR_LIGHT_BG)
    add_slide_header(slide5, "Technical deep-dive", "Speech-to-Text & Automatic Translation")
    
    # Left Box - Description & Code Snippets
    left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.8))
    tf5 = left_box.text_frame
    tf5.word_wrap = True
    
    p = tf5.paragraphs[0]
    p.text = "OpenAI Whisper Model Deployment"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    bullets = [
        ("Local Inference", "Utilizes the high-quality OpenAI Whisper Large-v2 model entirely on the local GPU/CPU. Eliminates cloud transcription costs."),
        ("Multilingual Translation", "Specifically engineered with `language='hi'` and `task='translate'` inside the python codebase. Spoken Hindi/Hinglish lectures are transcribed and translated instantly to English text segments."),
        ("Temporal Segment Mapping", "Unlike standard transcribers, Whisper preserves exact timestamps (start/end in seconds) for every statement, allowing micro-targeting of video sections.")
    ]
    for b_title, b_desc in bullets:
        pb_t = tf5.add_paragraph()
        pb_t.text = f"•  {b_title}"
        pb_t.font.name = FONT_BODY
        pb_t.font.size = Pt(15)
        pb_t.font.bold = True
        pb_t.font.color.rgb = COLOR_TEXT_DARK
        pb_t.space_after = Pt(3)
        
        pb_d = tf5.add_paragraph()
        pb_d.text = b_desc
        pb_d.font.name = FONT_BODY
        pb_d.font.size = Pt(13)
        pb_d.font.color.rgb = COLOR_TEXT_MUTED
        pb_d.space_after = Pt(12)

    # Right Box - Mock Code / Snippet UI Card
    right_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_NAVY
    right_card.line.fill.background()
    
    rtb = slide5.shapes.add_textbox(Inches(7.3), Inches(2.2), Inches(4.9), Inches(4.2))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "Whisper Integration Code"
    rp1.font.name = FONT_TITLE
    rp1.font.size = Pt(13)
    rp1.font.bold = True
    rp1.font.color.rgb = COLOR_CYAN
    rp1.space_after = Pt(14)
    
    code_lines = [
        "import whisper",
        "import json",
        "",
        "# Load heavy v2 model locally",
        "model = whisper.load_model('large-v2')",
        "",
        "# Transcribe and Translate Hindi to English",
        "result = model.transcribe(",
        "    audio = 'audios/lecture.mp3',",
        "    language = 'hi',",
        "    task = 'translate',",
        "    word_timestamps = False",
        ")",
        "",
        "# Map structured segments to JSON chunks",
        "for segment in result['segments']:",
        "    chunks.append({",
        "        'start': segment['start'],",
        "        'end': segment['end'],",
        "        'text': segment['text']",
        "    })"
    ]
    for line in code_lines:
        rp_code = rtf.add_paragraph()
        rp_code.text = line
        rp_code.font.name = "Consolas"
        rp_code.font.size = Pt(10)
        rp_code.font.color.rgb = COLOR_WHITE
        rp_code.space_after = Pt(0)

    # ----------------------------------------------------
    # Slide 6: Embedding & Indexing Component (Ollama)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6, COLOR_LIGHT_BG)
    add_slide_header(slide6, "Technical deep-dive", "Vector Embeddings & DataFrame Indexing")
    
    # Left Box - Description
    left_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.8))
    tf6 = left_box.text_frame
    tf6.word_wrap = True
    
    p = tf6.paragraphs[0]
    p.text = "Ollama Local Vector Indexing"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    bullets = [
        ("Ollama Local API Service", "Hosting models locally on `http://localhost:11434`. Restricts data sharing and saves network bandwidth, guaranteeing privacy."),
        ("BGE-M3 State-of-the-Art Embeddings", "Uses BAAI's `bge-m3` embedding model. Highly optimized for RAG operations and capable of processing multi-lingual inputs into 1024-dimensional dense vectors."),
        ("Pandas & Joblib Indexing", "Transcripts and vectors are stored in a unified Pandas DataFrame. The DataFrame is serialized to a local file `embeddings.joblib` using `joblib.dump()`, securing high-speed retrieval.")
    ]
    for b_title, b_desc in bullets:
        pb_t = tf6.add_paragraph()
        pb_t.text = f"•  {b_title}"
        pb_t.font.name = FONT_BODY
        pb_t.font.size = Pt(15)
        pb_t.font.bold = True
        pb_t.font.color.rgb = COLOR_TEXT_DARK
        pb_t.space_after = Pt(3)
        
        pb_d = tf6.add_paragraph()
        pb_d.text = b_desc
        pb_d.font.name = FONT_BODY
        pb_d.font.size = Pt(13)
        pb_d.font.color.rgb = COLOR_TEXT_MUTED
        pb_d.space_after = Pt(12)

    # Right Box - Mock Code / Snippet UI Card
    right_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_NAVY
    right_card.line.fill.background()
    
    rtb = slide6.shapes.add_textbox(Inches(7.3), Inches(2.2), Inches(4.9), Inches(4.2))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "Ollama Embedding Code"
    rp1.font.name = FONT_TITLE
    rp1.font.size = Pt(13)
    rp1.font.bold = True
    rp1.font.color.rgb = COLOR_CYAN
    rp1.space_after = Pt(14)
    
    code_lines = [
        "import requests",
        "import joblib",
        "import pandas as pd",
        "",
        "def create_embedding(text_list):",
        "    r = requests.post(",
        "        'http://localhost:11434/api/embed',",
        "        json={",
        "            'model': 'bge-m3',",
        "            'input': text_list",
        "        }",
        "    )",
        "    return r.json()['embeddings']",
        "",
        "# Processing JSON transcript files",
        "embeddings = create_embedding([c['text'] for c in chunks])",
        "for i, chunk in enumerate(chunks):",
        "    chunk['embedding'] = embeddings[i]",
        "",
        "df = pd.DataFrame.from_records(chunks)",
        "joblib.dump(df, 'embeddings.joblib')"
    ]
    for line in code_lines:
        rp_code = rtf.add_paragraph()
        rp_code.text = line
        rp_code.font.name = "Consolas"
        rp_code.font.size = Pt(10)
        rp_code.font.color.rgb = COLOR_WHITE
        rp_code.space_after = Pt(0)

    # ----------------------------------------------------
    # Slide 7: RAG Pipeline - Search & Llama 3.2
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7, COLOR_LIGHT_BG)
    add_slide_header(slide7, "Technical deep-dive", "Retrieval & Grounded LLM Q&A")
    
    # Left Box - Content
    left_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.8))
    tf7 = left_box.text_frame
    tf7.word_wrap = True
    
    p = tf7.paragraphs[0]
    p.text = "Mathematical Retrieval & Generation"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    bullets = [
        ("Cosine Similarity Retrieval", "Queries are encoded, then measured mathematically using scikit-learn's `cosine_similarity` against all stored chunks. Highly accurate, sub-millisecond retrieval."),
        ("Top 5 Context Extraction", "The system selects the highest scoring 5 segments, compiling them as raw JSON strings including video titles and timestamp intervals."),
        ("Llama 3.2 Inference", "Prompts Llama 3.2 locally using Ollama (`/api/generate`). Instructs the LLM to restrict responses specifically to the retrieved contexts, eliminating hallucinations.")
    ]
    for b_title, b_desc in bullets:
        pb_t = tf7.add_paragraph()
        pb_t.text = f"•  {b_title}"
        pb_t.font.name = FONT_BODY
        pb_t.font.size = Pt(15)
        pb_t.font.bold = True
        pb_t.font.color.rgb = COLOR_TEXT_DARK
        pb_t.space_after = Pt(3)
        
        pb_d = tf7.add_paragraph()
        pb_d.text = b_desc
        pb_d.font.name = FONT_BODY
        pb_d.font.size = Pt(13)
        pb_d.font.color.rgb = COLOR_TEXT_MUTED
        pb_d.space_after = Pt(12)

    # Right Box - Mock Code / Snippet UI Card
    right_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_NAVY
    right_card.line.fill.background()
    
    rtb = slide7.shapes.add_textbox(Inches(7.3), Inches(2.2), Inches(4.9), Inches(4.2))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "Similarity & Inference Code"
    rp1.font.name = FONT_TITLE
    rp1.font.size = Pt(13)
    rp1.font.bold = True
    rp1.font.color.rgb = COLOR_CYAN
    rp1.space_after = Pt(14)
    
    code_lines = [
        "from sklearn.metrics.pairwise import cosine_similarity",
        "import numpy as np",
        "",
        "# 1. Fetch query embedding",
        "q_embed = create_embedding([incoming_query])[0]",
        "",
        "# 2. Calculate Cosine Similarity",
        "similarities = cosine_similarity(",
        "    np.vstack(df['embedding']), [q_embed]",
        ").flatten()",
        "",
        "# 3. Extract Top-5 Indices",
        "max_idx = similarities.argsort()[::-1][:5]",
        "context_df = df.loc[max_idx]",
        "",
        "# 4. Structured RAG Prompt Setup",
        "prompt = f'''",
        "Here are video subtitle chunks:",
        "{context_df.to_json(orient='records')}",
        "---",
        "Answer user query: '{incoming_query}'",
        "Include video name and exact timestamps...",
        "'''"
    ]
    for line in code_lines:
        rp_code = rtf.add_paragraph()
        rp_code.text = line
        rp_code.font.name = "Consolas"
        rp_code.font.size = Pt(10)
        rp_code.font.color.rgb = COLOR_WHITE
        rp_code.space_after = Pt(0)

    # ----------------------------------------------------
    # Slide 8: Prompt Design & Hallucination Defense
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8, COLOR_LIGHT_BG)
    add_slide_header(slide8, "RAG Prompt Engineering", "Grounded System Prompting")
    
    # Left Content Column (Bullet points)
    left_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.8))
    tf8 = left_box.text_frame
    tf8.word_wrap = True
    
    p = tf8.paragraphs[0]
    p.text = "Securing LLM Responses"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    bullets = [
        ("Explicit Context Anchor", "The system prompt acts as a boundary, preventing the LLM from drawing on its pre-existing training data for general queries. Responses are tied strictly to the video metadata."),
        ("Persona Constraint", "The prompt guides the LLM to speak in a helpful, human-like voice suited for a university teaching assistant, while maintaining absolute factual precision."),
        ("Strict Boundary Handling", "Includes defensive instructions: if the user's question cannot be found within the retrieved video segments, the LLM will politely reject to answer, resolving hallucinations.")
    ]
    for b_title, b_desc in bullets:
        pb_t = tf8.add_paragraph()
        pb_t.text = f"•  {b_title}"
        pb_t.font.name = FONT_BODY
        pb_t.font.size = Pt(15)
        pb_t.font.bold = True
        pb_t.font.color.rgb = COLOR_TEXT_DARK
        pb_t.space_after = Pt(3)
        
        pb_d = tf8.add_paragraph()
        pb_d.text = b_desc
        pb_d.font.name = FONT_BODY
        pb_d.font.size = Pt(13)
        pb_d.font.color.rgb = COLOR_TEXT_MUTED
        pb_d.space_after = Pt(12)

    # Right Box - Actual Prompt Visual Card
    right_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_NAVY
    right_card.line.color.rgb = COLOR_CYAN
    right_card.line.width = Pt(1)
    
    rtb = slide8.shapes.add_textbox(Inches(7.3), Inches(2.2), Inches(4.9), Inches(4.2))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "ACTUAL CODE PROMPT TEMPLATE"
    rp1.font.name = FONT_TITLE
    rp1.font.size = Pt(13)
    rp1.font.bold = True
    rp1.font.color.rgb = COLOR_CYAN
    rp1.space_after = Pt(14)
    
    prompt_text = (
        "I am teaching web development in my Sigma web development course. "
        "Here are video subtitle chunks containing video title, video number, "
        "start time in seconds, end time in seconds, the text at that time:\n\n"
        "[CONTEXT CHUNKS JSON DATA]\n"
        "---------------------------------\n"
        "\"[USER QUERY]\"\n\n"
        "User asked this question related to the video chunks, you have to answer in a human way "
        "(dont mention the above format, its just for you) where and how much content is taught in which video "
        "(in which video and at what timestamp) and guide the user to go to that particular video. "
        "If user asks unrelated question, tell him that you can only answer questions related to the course."
    )
    
    rp_p = rtf.add_paragraph()
    rp_p.text = prompt_text
    rp_p.font.name = FONT_BODY
    rp_p.font.size = Pt(11.5)
    rp_p.font.color.rgb = COLOR_WHITE
    rp_p.space_after = Pt(0)

    # ----------------------------------------------------
    # Slide 9: Core Advantages of This Local Design
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9, COLOR_LIGHT_BG)
    add_slide_header(slide9, "System Evaluation", "Key Advantages & Architectural Strengths")
    
    advantages = [
        ("🔒 100% Local Data Privacy", "No data is shared with third-party LLM providers. Since transcription (Whisper), indexing (BGE-M3), and inference (Llama 3.2) run completely on local hardware, student questions and course data remain secure."),
        ("💰 Zero Operating Costs", "Eliminates the cost barriers of proprietary APIs (like OpenAI GPT-4 or Pinecone). Pushing updates or handling queries costs zero dollars, regardless of student volume."),
        ("🌐 Precise Multi-lingual Support", "Bilingual instruction (Hinglish/Hindi audio) is perfectly parsed and saved as clean English vectors, bridging the gap between local colloquial instruction and semantic English queries."),
        ("⚡ High Performance Indexing", "Utilizing a customized joblib-pickled Pandas DataFrame indexes vectors inside RAM, returning results within milliseconds without high-complexity hardware infrastructure.")
    ]
    
    for i, (adv_title, adv_desc) in enumerate(advantages):
        left_pos = Inches(0.8) if i % 2 == 0 else Inches(6.8)
        top_pos = Inches(2.0) if i < 2 else Inches(4.3)
        
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.7), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)
        
        tb = slide9.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), Inches(5.3), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p1 = tf.paragraphs[0]
        p1.text = adv_title
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_NAVY
        p1.space_after = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = adv_desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------
    # Slide 10: Future Scope & Conclusion
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide10, COLOR_NAVY)
    
    # Left accent block
    accent_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar10.fill.solid()
    accent_bar10.fill.fore_color.rgb = COLOR_CYAN
    accent_bar10.line.fill.background()
    
    left_box = slide10.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(6.0), Inches(5.0))
    tf10 = left_box.text_frame
    tf10.word_wrap = True
    
    p1 = tf10.paragraphs[0]
    p1.text = "FUTURE PATHS"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_CYAN
    p1.space_after = Pt(14)
    
    p2 = tf10.add_paragraph()
    p2.text = "Expanding the RAG Assistant"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_after = Pt(20)
    
    points = [
        ("🖥  Streamlit Web UI", "Integrate the Python pipeline with a clean, interactive chat interface where students can paste any YouTube playlist link and start chatting instantly."),
        ("👁  Multimodal Visual RAG", "Upgrade the search indexing from text transcriptions to visual screen frames, utilizing CLIP models to match questions to slides or code editor windows."),
        ("🔁 Continuous playlist ingestion", "Establish webhooks or scheduled cron scripts that automatically detect new YouTube course video uploads, downloading and indexing them dynamically.")
    ]
    for pt_title, pt_desc in points:
        pt1 = tf10.add_paragraph()
        pt1.text = pt_title
        pt1.font.name = FONT_BODY
        pt1.font.size = Pt(15)
        pt1.font.bold = True
        pt1.font.color.rgb = COLOR_CYAN
        pt1.space_after = Pt(3)
        
        pt2 = tf10.add_paragraph()
        pt2.text = pt_desc
        pt2.font.name = FONT_BODY
        pt2.font.size = Pt(12)
        pt2.font.color.rgb = COLOR_WHITE
        pt2.space_after = Pt(12)

    # Right Box - Thank you / Q&A Card
    right_box = slide10.shapes.add_textbox(Inches(7.8), Inches(2.2), Inches(4.5), Inches(4.0))
    tf_thank = right_box.text_frame
    tf_thank.word_wrap = True
    
    pt = tf_thank.paragraphs[0]
    pt.text = "Thank You!"
    pt.font.name = FONT_TITLE
    pt.font.size = Pt(48)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_WHITE
    pt.space_after = Pt(10)
    
    pq = tf_thank.add_paragraph()
    pq.text = "Open for Questions & Discussion"
    pq.font.name = FONT_BODY
    pq.font.size = Pt(18)
    pq.font.color.rgb = COLOR_CYAN
    pq.space_after = Pt(30)
    
    pc = tf_thank.add_paragraph()
    pc.text = "GitHub: PraDarch/Rag-Based-AI-Teaching-Assistance"
    pc.font.name = FONT_BODY
    pc.font.size = Pt(13)
    pc.font.color.rgb = COLOR_WHITE
    
    # Save the presentation
    prs.save("RAG_AI_Teaching_Assistant_Presentation.pptx")
    print("Presentation created successfully as RAG_AI_Teaching_Assistant_Presentation.pptx!")

if __name__ == "__main__":
    create_presentation()
